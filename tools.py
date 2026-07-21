"""
Tools for the Groq GAIA agent. Each tool returns a string and never raises into
the agent loop on expected failures (returns "ERROR: ..." instead). Outputs are
truncated to keep Groq's per-minute token budget small.

Coverage across GAIA question types:
  read_file             txt/md/py/json/csv/xlsx/pdf/docx/pptx  (local path OR URL)
  analyze_image         images: charts, chess, screenshots, OCR (Groq vision)
  transcribe_audio      mp3/wav/m4a/flac  (Groq Whisper)
  get_youtube_transcript spoken content of a YouTube video
  analyze_youtube_video  VISUAL content of a YouTube video (frame sampling + vision)
  visit_webpage         read a URL, with keyword windowing for long pages
plus smolagents' DuckDuckGoSearchTool and WikipediaSearchTool (added in agent.py).
"""

import os
import json
import base64
import mimetypes
import tempfile

from smolagents import tool

VISION_MODEL = os.getenv("GROQ_VISION_MODEL", os.getenv("GROQ_MODEL_ID", "groq/qwen/qwen3.6-27b"))
WHISPER_MODEL = os.getenv("GROQ_WHISPER_MODEL", "whisper-large-v3-turbo")

# Tight budget: smaller outputs mean fewer rate-limit stalls on the free tier.
_MAX_CHARS = 1200

# Optional proxy for YouTube (Colab/cloud IPs are usually blocked). Set e.g.
#   export YOUTUBE_PROXY="http://user:pass@host:port"
_YT_PROXY = os.getenv("YOUTUBE_PROXY", "").strip()


def _truncate(text: str, limit: int = _MAX_CHARS) -> str:
    if text is None:
        return ""
    if len(text) <= limit:
        return text
    return text[:limit] + f"\n\n...[truncated, {len(text) - limit} more chars]"


def _download(url: str) -> str:
    """Download a URL to a temp file, preserving a sensible extension."""
    import requests
    from urllib.parse import urlparse
    headers = {"User-Agent": "Mozilla/5.0 (compatible; GaiaAgent/1.0)"}
    r = requests.get(url, headers=headers, timeout=60)
    r.raise_for_status()
    ext = os.path.splitext(urlparse(url).path)[1]
    if not ext:
        ext = mimetypes.guess_extension(r.headers.get("content-type", "").split(";")[0]) or ""
    path = os.path.join(tempfile.gettempdir(), f"dl_{abs(hash(url)) % 10**8}{ext}")
    with open(path, "wb") as f:
        f.write(r.content)
    return path


# --------------------------------------------------------------------------- #
# File reading (local path OR http URL)                                       #
# --------------------------------------------------------------------------- #
@tool
def read_file(file_path: str) -> str:
    """Read a local file (or a file at an http/https URL) and return its text.

    Handles: .txt/.md/.py/.log/.xml/.html/.tsv (text), .json, .csv/.xlsx/.xls
    (tabular preview), .pdf (pdfplumber), .docx (python-docx), .pptx
    (python-pptx). For images use analyze_image; for audio use transcribe_audio.

    To COMPUTE over a full spreadsheet/CSV, load it yourself:
    `import pandas as pd; df = pd.read_excel(file_path)`.

    Args:
        file_path: Local path, or an http(s) URL to download and read.
    """
    if file_path.startswith(("http://", "https://")):
        try:
            file_path = _download(file_path)
        except Exception as e:
            return f"ERROR downloading {file_path}: {type(e).__name__}: {e}"

    if not os.path.exists(file_path):
        return f"ERROR: file not found at {file_path}"

    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext in {".txt", ".md", ".py", ".log", ".xml", ".html", ".tsv", ""}:
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return _truncate(f.read())

        if ext == ".json":
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return _truncate(json.dumps(json.load(f), indent=2, ensure_ascii=False))

        if ext == ".csv":
            import pandas as pd
            df = pd.read_csv(file_path)
            return _truncate(
                f"CSV shape {df.shape}; columns={list(df.columns)}\n"
                f"(compute over the full file with pd.read_csv('{file_path}'))\n\n"
                f"{df.head(30).to_string()}"
            )

        if ext in {".xlsx", ".xls"}:
            import pandas as pd
            sheets = pd.read_excel(file_path, sheet_name=None)
            out = [f"=== sheet: {n} (shape {d.shape}); columns={list(d.columns)} ===\n"
                   f"{d.head(30).to_string()}" for n, d in sheets.items()]
            return _truncate(
                f"(compute with pd.read_excel('{file_path}', sheet_name=None))\n\n"
                + "\n\n".join(out)
            )

        if ext == ".pdf":
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                pages = [p.extract_text() or "" for p in pdf.pages]
            return _truncate("\n\n".join(pages))

        if ext == ".docx":
            import docx
            d = docx.Document(file_path)
            parts = [p.text for p in d.paragraphs if p.text.strip()]
            for ti, tbl in enumerate(d.tables):
                rows = [" | ".join(c.text for c in row.cells) for row in tbl.rows]
                parts.append(f"[table {ti}]\n" + "\n".join(rows))
            return _truncate("\n".join(parts))

        if ext == ".pptx":
            import pptx
            prs = pptx.Presentation(file_path)
            parts = []
            for si, slide in enumerate(prs.slides, 1):
                texts = [sh.text for sh in slide.shapes
                         if getattr(sh, "has_text_frame", False) and sh.text.strip()]
                if texts:
                    parts.append(f"[slide {si}]\n" + "\n".join(texts))
            return _truncate("\n\n".join(parts))

        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return _truncate(f.read())

    except Exception as e:
        return f"ERROR reading {file_path}: {type(e).__name__}: {e}"


# --------------------------------------------------------------------------- #
# Vision (Groq)                                                               #
# --------------------------------------------------------------------------- #
def _data_uri(file_path: str) -> str:
    mime = mimetypes.guess_type(file_path)[0] or "image/png"
    with open(file_path, "rb") as f:
        b64 = base64.b64encode(f.read()).decode("utf-8")
    return f"data:{mime};base64,{b64}"


def _vision_call(image_paths: list[str], question: str) -> str:
    import litellm
    content = [{"type": "text", "text": question}]
    for p in image_paths:
        content.append({"type": "image_url", "image_url": {"url": _data_uri(p)}})
    resp = litellm.completion(
        model=VISION_MODEL, temperature=0.0,
        messages=[{"role": "user", "content": content}],
    )
    return resp.choices[0].message.content or ""


@tool
def analyze_image(file_path: str, question: str) -> str:
    """Answer a question about an image (photo, chart, diagram, screenshot,
    chessboard); handles OCR too. Use for any image file.

    Args:
        file_path: Local path (or http URL) to the image.
        question: The specific question to answer about the image.
    """
    if file_path.startswith(("http://", "https://")):
        try:
            file_path = _download(file_path)
        except Exception as e:
            return f"ERROR downloading image: {e}"
    if not os.path.exists(file_path):
        return f"ERROR: image not found at {file_path}"
    try:
        return _truncate(_vision_call([file_path], question))
    except Exception as e:
        try:
            from PIL import Image
            from pytesseract import image_to_string
            text = image_to_string(Image.open(file_path))
            if text.strip():
                return _truncate("Vision unavailable; OCR text only:\n" + text)
        except Exception:
            pass
        return f"ERROR analyzing image: {type(e).__name__}: {e}"


# --------------------------------------------------------------------------- #
# Audio (Groq Whisper)                                                        #
# --------------------------------------------------------------------------- #
@tool
def transcribe_audio(file_path: str) -> str:
    """Transcribe a local (or http URL) audio file (mp3/wav/m4a/flac) with Groq
    Whisper. Use for any audio attachment, e.g. a voice memo.

    Args:
        file_path: Local path or http URL to the audio file.
    """
    if file_path.startswith(("http://", "https://")):
        try:
            file_path = _download(file_path)
        except Exception as e:
            return f"ERROR downloading audio: {e}"
    if not os.path.exists(file_path):
        return f"ERROR: audio not found at {file_path}"
    try:
        import litellm
        with open(file_path, "rb") as f:
            resp = litellm.transcription(
                model=f"groq/{WHISPER_MODEL}", file=f,
                temperature=0.0, response_format="text",
            )
        text = getattr(resp, "text", None) or (resp if isinstance(resp, str) else str(resp))
        return _truncate(text)
    except Exception as e:
        return f"ERROR transcribing audio: {type(e).__name__}: {e}"


# --------------------------------------------------------------------------- #
# YouTube                                                                     #
# --------------------------------------------------------------------------- #
def _yt_id(url_or_id: str) -> str:
    vid = url_or_id.strip()
    if "watch?v=" in vid:
        vid = vid.split("watch?v=")[1].split("&")[0]
    elif "youtu.be/" in vid:
        vid = vid.split("youtu.be/")[1].split("?")[0]
    elif "shorts/" in vid:
        vid = vid.split("shorts/")[1].split("?")[0]
    return vid


@tool
def get_youtube_transcript(url_or_id: str) -> str:
    """Fetch the transcript (spoken words) of a YouTube video. Use when the
    answer is in what is SAID. Accepts a URL or a bare video id. Note: YouTube
    blocks cloud/Colab IPs; set YOUTUBE_PROXY (a residential proxy) if this
    returns RequestBlocked.

    Args:
        url_or_id: A YouTube URL or the raw video id.
    """
    try:
        from youtube_transcript_api import YouTubeTranscriptApi
    except Exception as e:
        return f"ERROR: youtube-transcript-api not installed: {e}"
    vid = _yt_id(url_or_id)
    try:
        # Configure a proxy if provided (helps get around cloud-IP blocks).
        api = None
        if _YT_PROXY:
            try:
                from youtube_transcript_api.proxies import GenericProxyConfig
                api = YouTubeTranscriptApi(
                    proxy_config=GenericProxyConfig(http_url=_YT_PROXY, https_url=_YT_PROXY)
                )
            except Exception:
                api = None
        api = api or YouTubeTranscriptApi()
        fetched = api.fetch(vid)
        try:
            text = " ".join(s["text"] for s in fetched.to_raw_data())
        except AttributeError:
            text = " ".join(sn.text for sn in fetched)
        return _truncate(text)
    except Exception as e:
        return f"ERROR fetching transcript for '{vid}': {type(e).__name__}: {e}"


@tool
def analyze_youtube_video(url: str, question: str, every_n_seconds: int = 15) -> str:
    """Answer a question about the VISUAL content of a YouTube video by sampling
    frames and sending them to the vision model. Use ONLY when the answer needs
    what is SEEN (objects, counts, colors, on-screen text). For spoken content
    use get_youtube_transcript instead. Needs yt-dlp + imageio; may fail if the
    Colab IP is blocked by YouTube.

    Args:
        url: The YouTube video URL.
        question: What to determine from the video's visuals.
        every_n_seconds: Sampling interval; larger = fewer frames.
    """
    try:
        import yt_dlp
        import imageio.v2 as imageio
        from PIL import Image
    except Exception as e:
        return f"ERROR: video deps missing (pip install yt-dlp imageio): {e}"
    try:
        with tempfile.TemporaryDirectory() as td:
            opts = {"format": "worstvideo[ext=mp4]/mp4", "quiet": True,
                    "outtmpl": os.path.join(td, "v.%(ext)s"), "force_ipv4": True,
                    "socket_timeout": 20, "retries": 1}
            if _YT_PROXY:
                opts["proxy"] = _YT_PROXY
            with yt_dlp.YoutubeDL(opts) as ydl:
                ydl.extract_info(url, download=True)
            vf = next((os.path.join(td, f) for f in os.listdir(td)
                       if f.endswith((".mp4", ".mkv", ".webm"))), None)
            if not vf:
                return "ERROR: could not download video"
            reader = imageio.get_reader(vf)
            fps = reader.get_meta_data().get("fps", 30) or 30
            interval = max(1, int(fps * every_n_seconds))
            frames, paths = [], []
            for i, fr in enumerate(reader):
                if i % interval == 0:
                    p = os.path.join(td, f"f{i}.jpg")
                    Image.fromarray(fr).save(p)
                    paths.append(p)
                if len(paths) >= 6:   # cap frames to protect the token budget
                    break
            if not paths:
                return "ERROR: no frames extracted"
            return _truncate(_vision_call(paths, question))
    except Exception as e:
        return f"ERROR analyzing video: {type(e).__name__}: {e}"


# --------------------------------------------------------------------------- #
# Web page reader (with keyword windowing)                                    #
# --------------------------------------------------------------------------- #
def _fetch_page_text(url: str) -> str:
    import re
    import requests
    headers = {"User-Agent": "Mozilla/5.0 (compatible; GaiaAgent/1.0)"}
    r = requests.get(url, headers=headers, timeout=30)
    r.raise_for_status()
    text = None
    try:
        import trafilatura
        text = trafilatura.extract(r.text, include_links=False, include_tables=True)
    except Exception:
        text = None
    if not text:
        try:
            from markdownify import markdownify as md
            text = md(r.text)
        except Exception:
            text = r.text
    return re.sub(r"\n{3,}", "\n\n", text or "").strip()


@tool
def visit_webpage(url: str, search: str = "") -> str:
    """Fetch a web page's main text (tables kept). Only a window is returned; if
    the part you need isn't shown, call again with search='<keyword>' to center
    the window on that keyword. Do NOT re-fetch with identical arguments.

    Args:
        url: The full URL to fetch.
        search: Optional keyword to jump to (case-insensitive).
    """
    try:
        text = _fetch_page_text(url)
    except Exception as e:
        return f"ERROR fetching {url}: {type(e).__name__}: {e}"
    if not text:
        return f"ERROR: no readable text extracted from {url}"

    if search:
        idx = text.lower().find(search.lower())
        if idx == -1:
            import re
            heads = re.findall(r"^#{1,6}\s*(.+)$|^\s*(.+)\n[=-]{3,}\s*$",
                               text, flags=re.MULTILINE)
            flat = [h[0] or h[1] for h in heads][:30]
            hint = "; ".join(flat) if flat else "no clear headings"
            return (f"'{search}' not found. Headings present: {hint}. "
                    f"Page is {len(text)} chars.")
        start = max(0, idx - 200)
        return (f"[window around '{search}', page {len(text)} chars]\n"
                + text[start:start + _MAX_CHARS])

    if len(text) > _MAX_CHARS:
        return (f"[first {_MAX_CHARS} of {len(text)} chars; call again with "
                f"search='<keyword>' to jump ahead]\n" + text[:_MAX_CHARS])
    return text
